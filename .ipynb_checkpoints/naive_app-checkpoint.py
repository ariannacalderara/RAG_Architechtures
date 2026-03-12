import chromadb
from chromadb.utils import embedding_functions
import requests
import fitz
import pytesseract
from PIL import Image
import io
import os
import docx
from pptx import Presentation
import pandas as pd
from datetime import datetime
import streamlit as st

# ── ReportLab imports for PDF export ──────────────────────────────────────────
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    HRFlowable, PageBreak, KeepTogether
)
from reportlab.lib.enums import TA_LEFT, TA_CENTER

TEMP_DIR = "temp_files"
os.makedirs(TEMP_DIR, exist_ok=True)

CHROMA_DIR = "chroma_data"
client = chromadb.PersistentClient(path=CHROMA_DIR)
embed_fn = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
collection = client.get_or_create_collection(name="course_docs", embedding_function=embed_fn)

OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL = "tinyllama"

# ── Test questions for Experiment 2 ───────────────────────────────────────────
TEST_QUESTIONS = [
    {"id": "Q01", "question": "Summarise the main argument that spans sections 2 and 3 of the lecture notes.", "category": "Chunking"},
    {"id": "Q02", "question": "What conclusion does the author draw at the end of the document after building up the argument in the earlier sections?", "category": "Chunking"},
    {"id": "Q03", "question": "What is the value in the third row of the main data table?", "category": "Table Parsing"},
    {"id": "Q04", "question": "Which category has the highest total across all sheets in the spreadsheet?", "category": "Cross-sheet Reasoning"},
    {"id": "Q05", "question": "What do the merged header cells in the table represent?", "category": "Table Parsing"},
    {"id": "Q06", "question": "What does the diagram on slide 4 illustrate?", "category": "Image Blind Spot"},
    {"id": "Q07", "question": "Describe the figure shown in the scanned PDF page.", "category": "OCR Quality"},
    {"id": "Q08", "question": "What exact percentage was mentioned for customer satisfaction in 2023?", "category": "Hallucination"},
    {"id": "Q09", "question": "Who wrote this document and when was it last updated?", "category": "Hallucination"},
    {"id": "Q10", "question": "What is written in the footnote at the bottom of page 2?", "category": "Layout Parsing"},
    {"id": "Q11", "question": "Explain the relationship between the two columns of text on the title page.", "category": "Layout Parsing"},
    {"id": "Q12", "question": "What are the learning objectives listed at the very beginning of chapter 1?", "category": "Retrieval Miss"},
]

FAILURE_CATEGORIES = [
    "R — Retrieval Miss",
    "C — Chunking Break",
    "T — Table / Format Loss",
    "I — Image Blind Spot",
    "H — Hallucination",
    "L — Layout Error",
    "O — Other",
]

ARCHITECTURES = ["Naive RAG"]

# ── Extractors ────────────────────────────────────────────────────────────────

def extract_text_from_pdf(file_path):
    chunks = []
    doc = fitz.open(file_path)
    for page in doc:
        page_text = page.get_text().strip()
        if page_text:
            chunks.append(page_text)
        else:
            pix = page.get_pixmap(dpi=250)
            img = Image.open(io.BytesIO(pix.pil_tobytes()))
            ocr_text = pytesseract.image_to_string(img).strip()
            if ocr_text:
                chunks.append(ocr_text)
    doc.close()
    return chunks

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    chunks, current_chunk = [], []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading") and current_chunk:
            chunks.append("\n".join(current_chunk))
            current_chunk = [text]
        else:
            current_chunk.append(text)
    if current_chunk:
        chunks.append("\n".join(current_chunk))
    return chunks if chunks else ["No text found in document."]

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if slide_text:
            chunks.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
    return chunks if chunks else ["No text found in presentation."]

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    chunk_size = 500
    return [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]

def extract_text_from_excel(file_path):
    chunks = []
    xl = pd.ExcelFile(file_path)
    for sheet_name in xl.sheet_names:
        df = pd.read_excel(xl, sheet_name=sheet_name).dropna(how="all").fillna("")
        chunks.append(f"[Sheet: {sheet_name}]\n" + df.to_string(index=False))
    return chunks if chunks else ["No data found in spreadsheet."]

def extract_chunks(file_path, file_name):
    ext = os.path.splitext(file_name)[1].lower()
    if ext == ".pdf":    return extract_text_from_pdf(file_path)
    elif ext == ".docx": return extract_text_from_docx(file_path)
    elif ext == ".pptx": return extract_text_from_pptx(file_path)
    elif ext == ".txt":  return extract_text_from_txt(file_path)
    elif ext in [".xlsx", ".xls"]: return extract_text_from_excel(file_path)
    return []

# ── RAG helpers ───────────────────────────────────────────────────────────────

def retrieve_chunks(question, n=3):
    results = collection.query(query_texts=[question], n_results=n)
    docs  = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    return docs, metas

def ask_llm(context, question):
    prompt = (
        "You are a helpful academic tutor for WU Vienna students. "
        "Use the following course material excerpts to answer the question.\n\n"
        f"Context:\n{context}\n\nQuestion: {question}\nAnswer:"
    )
    payload = {"model": MODEL, "prompt": prompt, "stream": False}
    try:
        res = requests.post(OLLAMA_URL, json=payload, timeout=120)
        return res.json().get("response", "").strip() if res.status_code == 200 else f"Error: {res.status_code}"
    except Exception as e:
        return f"Error communicating with Ollama: {e}"

# ── PDF export ────────────────────────────────────────────────────────────────

def build_failure_pdf(results: list, architecture: str) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm
    )
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle("CustomTitle", parent=styles["Title"],
                                 fontSize=18, spaceAfter=6, textColor=colors.HexColor("#1a1a2e"))
    subtitle_style = ParagraphStyle("Subtitle", parent=styles["Normal"],
                                    fontSize=11, spaceAfter=4, textColor=colors.HexColor("#444444"))
    h2_style = ParagraphStyle("H2", parent=styles["Heading2"],
                               fontSize=13, spaceBefore=10, spaceAfter=4,
                               textColor=colors.HexColor("#1a1a2e"))
    label_style = ParagraphStyle("Label", parent=styles["Normal"],
                                 fontSize=8, textColor=colors.HexColor("#888888"),
                                 spaceAfter=2)
    body_style = ParagraphStyle("Body", parent=styles["Normal"],
                                fontSize=10, leading=14, spaceAfter=6)
    chunk_style = ParagraphStyle("Chunk", parent=styles["Normal"],
                                 fontSize=9, leading=13, textColor=colors.HexColor("#333333"),
                                 backColor=colors.HexColor("#f5f5f5"),
                                 leftIndent=8, rightIndent=8, spaceAfter=4)
    answer_style = ParagraphStyle("Answer", parent=styles["Normal"],
                                  fontSize=10, leading=14,
                                  textColor=colors.HexColor("#1a1a2e"),
                                  leftIndent=8)

    story = []

    # ── Cover page ────────────────────────────────────────────────────────────
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph("Experiment 2: Failure Analysis", title_style))
    story.append(Paragraph("RAG Architecture Diagnosis Worksheet", subtitle_style))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#cccccc"), spaceAfter=12))
    story.append(Spacer(1, 0.5*cm))

    meta_data = [
        ["Architecture tested:", architecture],
        ["Date:", datetime.now().strftime("%d %B %Y")],
        ["Total cases:", str(len(results))],
    ]
    meta_table = Table(meta_data, colWidths=[5*cm, 10*cm])
    meta_table.setStyle(TableStyle([
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("TEXTCOLOR", (0, 0), (0, -1), colors.HexColor("#888888")),
        ("FONTNAME", (0, 0), (0, -1), "Helvetica"),
        ("FONTNAME", (1, 0), (1, -1), "Helvetica-Bold"),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(meta_table)
    story.append(Spacer(1, 1.5*cm))

    # ── Instructions ──────────────────────────────────────────────────────────
    story.append(Paragraph("Instructions", h2_style))
    instructions = (
        "For each case below, read the question, the retrieved context chunks, and the LLM answer. "
        "Then fill in the diagnosis box at the bottom of each case:<br/><br/>"
        "<b>1. Which RAG architecture most likely produced this failure?</b> "
        "Circle one: <i>Naive RAG &nbsp;|&nbsp; DeepDoc RAG &nbsp;|&nbsp; Tablebook RAG</i><br/><br/>"
        "<b>2. What type of failure is this?</b> Circle one code:<br/>"
        "&nbsp;&nbsp;<b>R</b> — Retrieval Miss &nbsp;&nbsp; "
        "<b>C</b> — Chunking Break &nbsp;&nbsp; "
        "<b>T</b> — Table/Format Loss &nbsp;&nbsp; "
        "<b>I</b> — Image Blind Spot &nbsp;&nbsp; "
        "<b>H</b> — Hallucination &nbsp;&nbsp; "
        "<b>L</b> — Layout Error<br/><br/>"
        "<b>3. Explain in 1–2 sentences why this failure occurred.</b>"
    )
    story.append(Paragraph(instructions, body_style))
    story.append(PageBreak())

    # ── One case per page ─────────────────────────────────────────────────────
    for i, r in enumerate(results):
        block = []

        # Header bar
        header_data = [[
            Paragraph(f"Case {r['id']}", ParagraphStyle("CaseNum", parent=styles["Normal"],
                      fontSize=13, textColor=colors.white, fontName="Helvetica-Bold")),
            Paragraph(r["category"], ParagraphStyle("CatLabel", parent=styles["Normal"],
                      fontSize=10, textColor=colors.HexColor("#ccddff")))
        ]]
        header_table = Table(header_data, colWidths=[3*cm, 13.5*cm])
        header_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#1a1a2e")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (0, 0), 10),
            ("LEFTPADDING", (1, 0), (1, 0), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ("ROUNDEDCORNERS", [4, 4, 0, 0]),
        ]))
        block.append(header_table)
        block.append(Spacer(1, 0.3*cm))

        # Question
        block.append(Paragraph("QUESTION", label_style))
        block.append(Paragraph(r["question"], body_style))
        block.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#dddddd"), spaceAfter=8))

        # Retrieved chunks
        block.append(Paragraph("RETRIEVED CONTEXT CHUNKS", label_style))
        for j, chunk in enumerate(r["chunks"], 1):
            source = r["sources"][j-1] if j-1 < len(r["sources"]) else "unknown"
            block.append(Paragraph(f"<b>Chunk {j}</b> &nbsp; <font color='#888888' size='8'>[{source}]</font>",
                                   ParagraphStyle("ChunkHead", parent=styles["Normal"], fontSize=9, spaceAfter=2)))
            preview = chunk[:600] + ("…" if len(chunk) > 600 else "")
            # Escape XML special chars
            preview = preview.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            block.append(Paragraph(preview, chunk_style))
        block.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#dddddd"), spaceAfter=8))

        # LLM answer
        block.append(Paragraph("LLM ANSWER", label_style))
        answer_text = r["answer"].replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        block.append(Paragraph(answer_text[:800] + ("…" if len(answer_text) > 800 else ""), answer_style))
        block.append(Spacer(1, 0.4*cm))
        block.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#1a1a2e"), spaceAfter=8))

        # Diagnosis box
        block.append(Paragraph("YOUR DIAGNOSIS", ParagraphStyle("DiagLabel", parent=styles["Normal"],
                               fontSize=9, textColor=colors.HexColor("#1a1a2e"),
                               fontName="Helvetica-Bold", spaceAfter=6)))

        diag_data = [
            [Paragraph("<b>Architecture:</b>  Naive RAG &nbsp;&nbsp; | &nbsp;&nbsp; DeepDoc RAG &nbsp;&nbsp; | &nbsp;&nbsp; Tablebook RAG",
                       ParagraphStyle("DiagText", parent=styles["Normal"], fontSize=10))],
            [Paragraph("<b>Failure type:</b>  R &nbsp; C &nbsp; T &nbsp; I &nbsp; H &nbsp; L &nbsp; O",
                       ParagraphStyle("DiagText", parent=styles["Normal"], fontSize=10))],
            [Paragraph("<b>Explanation:</b><br/><br/><br/>",
                       ParagraphStyle("DiagText", parent=styles["Normal"], fontSize=10))],
        ]
        diag_table = Table(diag_data, colWidths=[16.5*cm])
        diag_table.setStyle(TableStyle([
            ("BOX", (0, 0), (-1, -1), 1, colors.HexColor("#1a1a2e")),
            ("INNERGRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#f9f9ff")),
            ("TOPPADDING", (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ("LEFTPADDING", (0, 0), (-1, -1), 10),
        ]))
        block.append(diag_table)

        story.append(KeepTogether(block))
        if i < len(results) - 1:
            story.append(PageBreak())

    doc.build(story)
    return buf.getvalue()

# ── Streamlit UI ──────────────────────────────────────────────────────────────

st.set_page_config(page_title="WU Vienna AI Course Tutor", page_icon="🎓", layout="wide")
st.title("🎓💬 WU Vienna AI Course Tutor (Naive RAG)")
st.write("Upload course materials and ask questions. All processing is local and private!")

# ── Tab layout ────────────────────────────────────────────────────────────────
tab_chat, tab_export = st.tabs(["💬 Chat", "📋 Experiment 2 — Failure Export"])

# ── Chat tab ──────────────────────────────────────────────────────────────────
with tab_chat:
    uploaded_files = st.file_uploader(
        "Upload course material(s)",
        type=["pdf", "docx", "pptx", "txt", "xlsx", "xls"],
        accept_multiple_files=True
    )

    if st.button("Ingest Documents") and uploaded_files:
        for file in uploaded_files:
            file_path = os.path.join(TEMP_DIR, file.name)
            with open(file_path, "wb") as f:
                f.write(file.getbuffer())
            chunks = extract_chunks(file_path, file.name)
            if not chunks:
                st.warning(f"⚠️ Could not extract text from {file.name}.")
                continue
            for idx, chunk in enumerate(chunks):
                collection.add(
                    documents=[chunk],
                    metadatas=[{"source": file.name, "chunk": idx+1}],
                    ids=[f"{file.name}_chunk{idx+1}"]
                )
        st.success(f"✅ Indexed {len(uploaded_files)} document(s) into the database.")

    query = st.text_input("Ask a question about your course:")
    if st.button("Get Answer") and query:
        chunks, metas = retrieve_chunks(query)
        context = "\n".join(chunks)
        answer = ask_llm(context, query)
        st.subheader("Answer:")
        st.write(answer)
        with st.expander("Show retrieved context"):
            for chunk in chunks:
                st.write(f"- {chunk[:1000]}{'...' if len(chunk) > 1000 else ''}")

# ── Export tab ────────────────────────────────────────────────────────────────
with tab_export:
    st.subheader("Generate Student Worksheet")
    st.write(
        "This runs all 12 test questions against the current RAG collection and exports "
        "a printable PDF worksheet — one failure case per page — for students to diagnose."
    )

    architecture_label = st.selectbox(
        "Which architecture is loaded in this instance?",
        ARCHITECTURES,
        help="This label appears on the worksheet header so students know which system produced the failures."
    )

    if st.button("🚀 Run all test questions & generate PDF"):
        results = []
        progress = st.progress(0)
        status = st.empty()

        for i, q in enumerate(TEST_QUESTIONS):
            status.write(f"Running {q['id']}: {q['question'][:60]}…")
            chunks, metas = retrieve_chunks(q["question"], n=3)
            context = "\n---\n".join(chunks) if chunks else ""
            answer = ask_llm(context, q["question"])
            sources = [m.get("source", "?") for m in metas]
            results.append({
                "id": q["id"],
                "question": q["question"],
                "category": q["category"],
                "chunks": chunks,
                "sources": sources,
                "answer": answer,
            })
            progress.progress((i + 1) / len(TEST_QUESTIONS))

        status.write("✅ All questions run — building PDF…")

        try:
            pdf_bytes = build_failure_pdf(results, architecture_label)
            filename = f"failure_worksheet_{architecture_label.lower().replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pdf"
            st.success(f"PDF ready! {len(results)} cases exported.")
            st.download_button(
                label="⬇️ Download Worksheet PDF",
                data=pdf_bytes,
                file_name=filename,
                mime="application/pdf",
            )
        except Exception as e:
            st.error(f"PDF generation failed: {e}")

    st.markdown("---")
    st.markdown("**Failure category codes (print and share with students):**")
    cols = st.columns(3)
    codes = [
        ("R", "Retrieval Miss", "Wrong chunks fetched"),
        ("C", "Chunking Break", "Answer split across chunk boundary"),
        ("T", "Table/Format Loss", "Table structure destroyed in parsing"),
        ("I", "Image Blind Spot", "Content was a diagram or image"),
        ("H", "Hallucination", "LLM invented facts not in context"),
        ("L", "Layout Error", "Multi-column or footnote parsed wrong"),
    ]
    for i, (code, name, desc) in enumerate(codes):
        with cols[i % 3]:
            st.markdown(f"**{code} — {name}**  \n{desc}")