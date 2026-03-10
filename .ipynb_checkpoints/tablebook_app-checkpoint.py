import streamlit as st
import chromadb
from chromadb.utils import embedding_functions
import requests
import fitz
import pytesseract
from PIL import Image
import io
import os
import docx
from pptx import Presentation
import pandas as pd
import pdfplumber

TEMP_DIR = "temp_files"
os.makedirs(TEMP_DIR, exist_ok=True)

CHROMA_DIR = "chroma_data"
client = chromadb.PersistentClient(path=CHROMA_DIR)
embed_fn = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
collection = client.get_or_create_collection(name="course_docs_tablebook", embedding_function=embed_fn)

# ── Extractors ────────────────────────────────────────────────────────────────

def extract_from_pdf(file_path):
    chunks = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            # Extract prose
            text = page.extract_text()
            if text and text.strip():
                chunks.append({"text": text.strip(), "type": "text", "page": i+1})
            # Extract tables separately
            for table in page.extract_tables():
                rows = [" | ".join([c or "" for c in row]) for row in table]
                table_text = "\n".join(rows)
                if table_text.strip():
                    chunks.append({"text": table_text, "type": "table", "page": i+1})
    return chunks

def extract_from_docx(file_path):
    doc = docx.Document(file_path)
    chunks = []
    current_chunk = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith("Heading") and current_chunk:
            chunks.append({"text": "\n".join(current_chunk), "type": "text", "page": 1})
            current_chunk = [text]
        else:
            current_chunk.append(text)
    if current_chunk:
        chunks.append({"text": "\n".join(current_chunk), "type": "text", "page": 1})
    # Extract tables from docx
    for table in doc.tables:
        rows = [" | ".join([cell.text.strip() for cell in row.cells]) for row in table.rows]
        table_text = "\n".join(rows)
        if table_text.strip():
            chunks.append({"text": table_text, "type": "table", "page": 1})
    return chunks if chunks else [{"text": "No text found.", "type": "text", "page": 1}]

def extract_from_pptx(file_path):
    prs = Presentation(file_path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            # Extract tables from slides
            if shape.has_table:
                rows = [" | ".join([cell.text.strip() for cell in row.cells]) for row in shape.table.rows]
                table_text = "\n".join(rows)
                if table_text.strip():
                    chunks.append({"text": f"[Slide {i+1} Table]\n{table_text}", "type": "table", "page": i+1})
        if slide_text:
            chunks.append({"text": f"[Slide {i+1}]\n" + "\n".join(slide_text), "type": "text", "page": i+1})
    return chunks if chunks else [{"text": "No text found.", "type": "text", "page": 1}]

def extract_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    chunk_size = 500
    return [{"text": content[i:i+chunk_size], "type": "text", "page": 1}
            for i in range(0, len(content), chunk_size)]

def extract_from_excel(file_path):
    chunks = []
    xl = pd.ExcelFile(file_path)
    for sheet_name in xl.sheet_names:
        df = pd.read_excel(xl, sheet_name=sheet_name)
        df = df.dropna(how="all").fillna("")
        sheet_text = f"[Sheet: {sheet_name}]\n" + df.to_string(index=False)
        chunks.append({"text": sheet_text, "type": "table", "page": 1})
    return chunks if chunks else [{"text": "No data found.", "type": "table", "page": 1}]

def extract_chunks(file_path, file_name):
    ext = os.path.splitext(file_name)[1].lower()
    if ext == ".pdf":
        return extract_from_pdf(file_path)
    elif ext == ".docx":
        return extract_from_docx(file_path)
    elif ext == ".pptx":
        return extract_from_pptx(file_path)
    elif ext == ".txt":
        return extract_from_txt(file_path)
    elif ext in [".xlsx", ".xls"]:
        return extract_from_excel(file_path)
    else:
        return []

# ── UI ────────────────────────────────────────────────────────────────────────

st.title("🎓📊 WU Vienna AI Course Tutor (TableBook RAG)")
st.write("Upload course materials and ask questions. All processing is local and private!")

uploaded_files = st.file_uploader(
    "Upload course material(s)",
    type=["pdf", "docx", "pptx", "txt", "xlsx", "xls"],
    accept_multiple_files=True
)

if st.button("Ingest Documents") and uploaded_files:
    for file in uploaded_files:
        file_path = os.path.join(TEMP_DIR, file.name)
        with open(file_path, "wb") as f:
            f.write(file.getbuffer())

        chunks = extract_chunks(file_path, file.name)
        if not chunks:
            st.warning(f"⚠️ Could not extract text from {file.name} — unsupported format.")
            continue

        for idx, chunk in enumerate(chunks):
            doc_id = f"{file.name}_chunk{idx+1}"
            metadata = {
                "source": file.name,
                "chunk": idx+1,
                "type": chunk["type"],    # "text" or "table"
                "page": chunk["page"]
            }
            collection.add(documents=[chunk["text"]], metadatas=[metadata], ids=[doc_id])

    st.success(f"✅ Indexed {len(uploaded_files)} document(s) into the database.")

# ── Query ─────────────────────────────────────────────────────────────────────

query = st.text_input("Ask a question about your course:")

if st.button("Get Answer") and query:
    # Dual retrieval — text and tables separately
    text_results = collection.query(
        query_texts=[query],
        n_results=2,
        where={"type": "text"}
    )
    table_results = collection.query(
        query_texts=[query],
        n_results=2,
        where={"type": "table"}
    )

    text_chunks = text_results["documents"][0] if text_results.get("documents") else []
    table_chunks = table_results["documents"][0] if table_results.get("documents") else []
    top_chunks = text_chunks + table_chunks

    context = "\n".join(top_chunks)[:800]

    ollama_url = "http://localhost:11434/api/generate"
    model = "tinyllama"

    prompt_text = (
        f"You are a helpful academic tutor for WU Vienna students. "
        f"Use the following course material excerpts to answer the question. "
        f"Some excerpts may be tables — use them carefully.\n\n"
        f"Context:\n{context}\n\nQuestion: {query}\nAnswer:"
    )

    payload = {"model": model, "prompt": prompt_text, "stream": False}

    try:
        res = requests.post(ollama_url, json=payload, timeout=1200)
        answer = res.json().get("response", "").strip() if res.status_code == 200 else f"Error: {res.status_code}"
    except Exception as e:
        answer = f"Error communicating with Ollama: {e}"

    st.subheader("Answer:")
    st.write(answer)

    with st.expander("Show retrieved context"):
        st.write("**📝 Text chunks:**")
        for chunk in text_chunks:
            st.write(f"- {chunk[:1000]}{'...' if len(chunk) > 1000 else ''}")
        st.write("**📊 Table chunks:**")
        for chunk in table_chunks:
            st.write(f"- {chunk[:1000]}{'...' if len(chunk) > 1000 else ''}")