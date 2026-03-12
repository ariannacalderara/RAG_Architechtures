import streamlit as st
import chromadb
from chromadb.utils import embedding_functions
import requests
import fitz
import pytesseract
from PIL import Image
import io
import os
import docx
from pptx import Presentation
import pandas as pd
import pdfplumber
from datetime import datetime

# ── ReportLab imports for PDF export ──────────────────────────────────────────
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table as RLTable, TableStyle,
    HRFlowable, PageBreak, KeepTogether
)

TEMP_DIR = "temp_files"
os.makedirs(TEMP_DIR, exist_ok=True)

CHROMA_DIR = "chroma_data"
client = chromadb.PersistentClient(path=CHROMA_DIR)
embed_fn = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
collection = client.get_or_create_collection(name="course_docs_tablebook", embedding_function=embed_fn)

OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL = "tinyllama"

# ── Test questions ─────────────────────────────────────────────────────────────
TEST_QUESTIONS = [
    {"id": "Q01", "question": "Summarise the main argument that spans sections 2 and 3 of the lecture notes.", "category": "Chunking"},
    {"id": "Q02", "question": "What conclusion does the author draw at the end of the document after building up the argument in the earlier sections?", "category": "Chunking"},
    {"id": "Q03", "question": "What is the value in the third row of the main data table?", "category": "Table Parsing"},
    {"id": "Q04", "question": "Which category has the highest total across all sheets in the spreadsheet?", "category": "Cross-sheet Reasoning"},
    {"id": "Q05", "question": "What do the merged header cells in the table represent?", "category": "Table Parsing"},
    {"id": "Q06", "question": "What does the diagram on slide 4 illustrate?", "category": "Image Blind Spot"},
    {"id": "Q07", "question": "Describe the figure shown in the scanned PDF page.", "category": "OCR Quality"},
    {"id": "Q08", "question": "What exact percentage was mentioned for customer satisfaction in 2023?", "category": "Hallucination"},
    {"id": "Q09", "question": "Who wrote this document and when was it last updated?", "category": "Hallucination"},
    {"id": "Q10", "question": "What is written in the footnote at the bottom of page 2?", "category": "Layout Parsing"},
    {"id": "Q11", "question": "Explain the relationship between the two columns of text on the title page.", "category": "Layout Parsing"},
    {"id": "Q12", "question": "What are the learning objectives listed at the very beginning of chapter 1?", "category": "Retrieval Miss"},
]

# ── Extractors ────────────────────────────────────────────────────────────────

def extract_from_pdf(file_path):
    chunks = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                chunks.append({"text": text.strip(), "type": "text", "page": i+1})
            for table in page.extract_tables():
                rows = [" | ".join([c or "" for c in row]) for row in table]
                table_text = "\n".join(rows)
                if table_text.strip():
                    chunks.append({"text": table_text, "type": "table", "page": i+1})
    return chunks

def extract_from_docx(file_path):
    doc = docx.Document(file_path)
    chunks = []
    current_chunk = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading") and current_chunk:
            chunks.append({"text": "\n".join(current_chunk), "type": "text", "page": 1})
            current_chunk = [text]
        else:
            current_chunk.append(text)
    if current_chunk:
        chunks.append({"text": "\n".join(current_chunk), "type": "text", "page": 1})
    for table in doc.tables:
        rows = [" | ".join([cell.text.strip() for cell in row.cells]) for row in table.rows]
        table_text = "\n".join(rows)
        if table_text.strip():
            chunks.append({"text": table_text, "type": "table", "page": 1})
    return chunks if chunks else [{"text": "No text found.", "type": "text", "page": 1}]

def extract_from_pptx(file_path):
    prs = Presentation(file_path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            if shape.has_table:
                rows = [" | ".join([cell.text.strip() for cell in row.cells]) for row in shape.table.rows]
                table_text = "\n".join(rows)
                if table_text.strip():
                    chunks.append({"text": f"[Slide {i+1} Table]\n{table_text}", "type": "table", "page": i+1})
        if slide_text:
            chunks.append({"text": f"[Slide {i+1}]\n" + "\n".join(slide_text), "type": "text", "page": i+1})
    return chunks if chunks else [{"text": "No text found.", "type": "text", "page": 1}]

def extract_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    chunk_size = 500
    return [{"text": content[i:i+chunk_size], "type": "text", "page": 1}
            for i in range(0, len(content), chunk_size)]

def extract_from_excel(file_path):
    chunks = []
    xl = pd.ExcelFile(file_path)
    for sheet_name in xl.sheet_names:
        df = pd.read_excel(xl, sheet_name=sheet_name).dropna(how="all").fillna("")
        sheet_text = f"[Sheet: {sheet_name}]\n" + df.to_string(index=False)
        chunks.append({"text": sheet_text, "type": "table", "page": 1})
    return chunks if chunks else [{"text": "No data found.", "type": "table", "page": 1}]

def extract_chunks(file_path, file_name):
    ext = os.path.splitext(file_name)[1].lower()
    if ext == ".pdf":         return extract_from_pdf(file_path)
    elif ext == ".docx":      return extract_from_docx(file_path)
    elif ext == ".pptx":      return extract_from_pptx(file_path)
    elif ext == ".txt":       return extract_from_txt(file_path)
    elif ext in [".xlsx", ".xls"]: return extract_from_excel(file_path)
    return []

# ── RAG helpers ───────────────────────────────────────────────────────────────

def retrieve_chunks(question, n=3):
    text_results = collection.query(
        query_texts=[question], n_results=2, where={"type": "text"})
    table_results = collection.query(
        query_texts=[question], n_results=2, where={"type": "table"})
    text_chunks  = text_results.get("documents",  [[]])[0]
    table_chunks = table_results.get("documents", [[]])[0]
    text_metas   = text_results.get("metadatas",  [[]])[0]
    table_metas  = table_results.get("metadatas", [[]])[0]
    return text_chunks + table_chunks, text_metas + table_metas

def ask_llm(context, question):
    prompt = (
        "You are a helpful academic tutor for WU Vienna students. "
        "Use the following course material excerpts to answer the question. "
        "Some excerpts may be tables — use them carefully.\n\n"
        f"Context:\n{context}\n\nQuestion: {question}\nAnswer:"
    )
    payload = {"model": MODEL, "prompt": prompt, "stream": False}
    try:
        res = requests.post(OLLAMA_URL, json=payload, timeout=120)
        return res.json().get("response", "").strip() if res.status_code == 200 else f"Error: {res.status_code}"
    except Exception as e:
        return f"Error communicating with Ollama: {e}"

# ── PDF export — amber palette ────────────────────────────────────────────────

def build_failure_pdf(results: list, architecture: str) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()

    HEADER_COLOR  = colors.HexColor("#4a2e00")   # dark amber/brown
    CHUNK_BG      = colors.HexColor("#fdf6ec")   # warm cream
    DIAG_BG       = colors.HexColor("#fdf6ec")
    DIVIDER_COLOR = colors.HexColor("#d4a96a")   # amber
    CHUNK_LABEL   = colors.HexColor("#d4a96a")

    title_style   = ParagraphStyle("CT", parent=styles["Title"],
                                   fontSize=18, spaceAfter=6,
                                   textColor=HEADER_COLOR)
    subtitle_style = ParagraphStyle("CS", parent=styles["Normal"],
                                    fontSize=11, spaceAfter=4,
                                    textColor=colors.HexColor("#7a5c2e"))
    h2_style      = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=13, spaceBefore=10, spaceAfter=4,
                                   textColor=HEADER_COLOR)
    label_style   = ParagraphStyle("LB", parent=styles["Normal"],
                                   fontSize=8, textColor=colors.HexColor("#888888"),
                                   spaceAfter=2)
    body_style    = ParagraphStyle("BD", parent=styles["Normal"],
                                   fontSize=10, leading=14, spaceAfter=6)
    chunk_style   = ParagraphStyle("CK", parent=styles["Normal"],
                                   fontSize=9, leading=13,
                                   textColor=colors.HexColor("#333333"),
                                   backColor=CHUNK_BG,
                                   leftIndent=8, rightIndent=8, spaceAfter=4)
    answer_style  = ParagraphStyle("AN", parent=styles["Normal"],
                                   fontSize=10, leading=14,
                                   textColor=HEADER_COLOR, leftIndent=8)

    story = []

    # Cover
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph("Experiment 2: Failure Analysis", title_style))
    story.append(Paragraph("RAG Architecture Diagnosis Worksheet", subtitle_style))
    story.append(HRFlowable(width="100%", thickness=1,
                            color=DIVIDER_COLOR, spaceAfter=12))
    story.append(Spacer(1, 0.5*cm))

    meta_data = [
        ["Architecture tested:", architecture],
        ["Date:", datetime.now().strftime("%d %B %Y")],
        ["Total cases:", str(len(results))],
    ]
    meta_tbl = RLTable(meta_data, colWidths=[5*cm, 10*cm])
    meta_tbl.setStyle(TableStyle([
        ("FONTSIZE",  (0,0), (-1,-1), 10),
        ("TEXTCOLOR", (0,0), (0,-1),  colors.HexColor("#888888")),
        ("FONTNAME",  (0,0), (0,-1),  "Helvetica"),
        ("FONTNAME",  (1,0), (1,-1),  "Helvetica-Bold"),
        ("BOTTOMPADDING", (0,0), (-1,-1), 6),
    ]))
    story.append(meta_tbl)
    story.append(Spacer(1, 1.5*cm))

    # Instructions
    story.append(Paragraph("Instructions", h2_style))
    story.append(Paragraph(
        "For each case below, read the question, the retrieved context chunks, and the LLM answer. "
        "Then fill in the diagnosis box:<br/><br/>"
        "<b>1. Which RAG architecture most likely produced this failure?</b> "
        "Circle one: <i>Naive RAG &nbsp;|&nbsp; DeepDoc RAG &nbsp;|&nbsp; Tablebook RAG</i><br/><br/>"
        "<b>2. What type of failure is this?</b> Circle one: "
        "<b>R</b> Retrieval Miss &nbsp; <b>C</b> Chunking Break &nbsp; "
        "<b>T</b> Table/Format Loss &nbsp; <b>I</b> Image Blind Spot &nbsp; "
        "<b>H</b> Hallucination &nbsp; <b>L</b> Layout Error<br/><br/>"
        "<b>3. Explain in 1–2 sentences why this failure occurred.</b>",
        body_style))
    story.append(PageBreak())

    # One case per page
    for i, r in enumerate(results):
        block = []

        header_data = [[
            Paragraph(f"Case {r['id']}",
                      ParagraphStyle("CN", parent=styles["Normal"],
                                     fontSize=13, textColor=colors.white,
                                     fontName="Helvetica-Bold")),
            Paragraph(r["category"],
                      ParagraphStyle("CC", parent=styles["Normal"],
                                     fontSize=10, textColor=CHUNK_LABEL))
        ]]
        header_tbl = RLTable(header_data, colWidths=[3*cm, 13.5*cm])
        header_tbl.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,-1), HEADER_COLOR),
            ("VALIGN",        (0,0), (-1,-1), "MIDDLE"),
            ("LEFTPADDING",   (0,0), (0,0),   10),
            ("LEFTPADDING",   (1,0), (1,0),   6),
            ("TOPPADDING",    (0,0), (-1,-1), 8),
            ("BOTTOMPADDING", (0,0), (-1,-1), 8),
        ]))
        block.append(header_tbl)
        block.append(Spacer(1, 0.3*cm))

        block.append(Paragraph("QUESTION", label_style))
        block.append(Paragraph(r["question"], body_style))
        block.append(HRFlowable(width="100%", thickness=0.5,
                                color=DIVIDER_COLOR, spaceAfter=8))
        #to delete, too much leading
        block.append(Paragraph("RETRIEVED CONTEXT CHUNKS", label_style))
        for j, chunk in enumerate(r["chunks"], 1):
            source = r["sources"][j-1] if j-1 < len(r["sources"]) else "unknown"
            block.append(Paragraph(
                f"<b>Chunk {j}</b> &nbsp; "
                f"<font color='#888888' size='8'>[{source}]</font>",
                ParagraphStyle("CH", parent=styles["Normal"], fontSize=9, spaceAfter=2)))
            preview = chunk[:600] + ("…" if len(chunk) > 600 else "")
            preview = preview.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            block.append(Paragraph(preview, chunk_style))

        block.append(HRFlowable(width="100%", thickness=0.5,
                                color=DIVIDER_COLOR, spaceAfter=8))

        block.append(Paragraph("LLM ANSWER", label_style))
        answer_text = r["answer"].replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
        block.append(Paragraph(
            answer_text[:800] + ("…" if len(answer_text) > 800 else ""),
            answer_style))
        block.append(Spacer(1, 0.4*cm))
        block.append(HRFlowable(width="100%", thickness=1,
                                color=HEADER_COLOR, spaceAfter=8))

        block.append(Paragraph("YOUR DIAGNOSIS",
                               ParagraphStyle("DL", parent=styles["Normal"],
                                              fontSize=9, fontName="Helvetica-Bold",
                                              textColor=HEADER_COLOR, spaceAfter=6)))
        diag_data = [
            [Paragraph("<b>Architecture:</b>  Naive RAG &nbsp;&nbsp; | &nbsp;&nbsp; "
                       "DeepDoc RAG &nbsp;&nbsp; | &nbsp;&nbsp; Tablebook RAG",
                       ParagraphStyle("DT", parent=styles["Normal"], fontSize=10))],
            [Paragraph("<b>Failure type:</b>  R &nbsp; C &nbsp; T &nbsp; I &nbsp; H &nbsp; L &nbsp; O",
                       ParagraphStyle("DT2", parent=styles["Normal"], fontSize=10))],
            [Paragraph("<b>Explanation:</b><br/><br/><br/>",
                       ParagraphStyle("DT3", parent=styles["Normal"], fontSize=10))],
        ]
        diag_tbl = RLTable(diag_data, colWidths=[16.5*cm])
        diag_tbl.setStyle(TableStyle([
            ("BOX",           (0,0), (-1,-1), 1,   HEADER_COLOR),
            ("INNERGRID",     (0,0), (-1,-1), 0.5, DIVIDER_COLOR),
            ("BACKGROUND",    (0,0), (-1,-1), DIAG_BG),
            ("TOPPADDING",    (0,0), (-1,-1), 8),
            ("BOTTOMPADDING", (0,0), (-1,-1), 8),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
        ]))
        block.append(diag_tbl)

        story.append(KeepTogether(block))
        if i < len(results) - 1:
            story.append(PageBreak())

    doc.build(story)
    return buf.getvalue()

# ── Streamlit UI ───────────────────────────────────────────────────────────────

st.set_page_config(page_title="WU Vienna AI Course Tutor", page_icon="🎓", layout="wide")
st.title("🎓📊 WU Vienna AI Course Tutor (TableBook RAG)")
st.write("Upload course materials and ask questions. All processing is local and private!")

tab_chat, tab_export = st.tabs(["💬 Chat", "📋 Experiment 2 — Failure Export"])

# ── Chat tab ───────────────────────────────────────────────────────────────────
with tab_chat:
    uploaded_files = st.file_uploader(
        "Upload course material(s)",
        type=["pdf", "docx", "pptx", "txt", "xlsx", "xls"],
        accept_multiple_files=True
    )

    if st.button("Ingest Documents") and uploaded_files:
        for file in uploaded_files:
            file_path = os.path.join(TEMP_DIR, file.name)
            with open(file_path, "wb") as f:
                f.write(file.getbuffer())
            chunks = extract_chunks(file_path, file.name)
            if not chunks:
                st.warning(f"⚠️ Could not extract text from {file.name}.")
                continue
            for idx, chunk in enumerate(chunks):
                collection.add(
                    documents=[chunk["text"]],
                    metadatas=[{"source": file.name, "chunk": idx+1,
                               "type": chunk["type"], "page": chunk["page"]}],
                    ids=[f"{file.name}_chunk{idx+1}"]
                )
        st.success(f"✅ Indexed {len(uploaded_files)} document(s) into the database.")

    query = st.text_input("Ask a question about your course:")
    if st.button("Get Answer") and query:
        chunks, metas = retrieve_chunks(query)
        context = "\n".join(chunks)[:800]
        answer = ask_llm(context, query)
        st.subheader("Answer:")
        st.write(answer)
        with st.expander("Show retrieved context"):
            text_chunks  = [c for c, m in zip(chunks, metas) if m.get("type") == "text"]
            table_chunks = [c for c, m in zip(chunks, metas) if m.get("type") == "table"]
            st.write("**📝 Text chunks:**")
            for chunk in text_chunks:
                st.write(f"- {chunk[:1000]}{'...' if len(chunk) > 1000 else ''}")
            st.write("**📊 Table chunks:**")
            for chunk in table_chunks:
                st.write(f"- {chunk[:1000]}{'...' if len(chunk) > 1000 else ''}")

# ── Export tab ─────────────────────────────────────────────────────────────────
with tab_export:
    st.subheader("Generate Student Worksheet")
    st.write(
        "Runs all 12 test questions against the current Tablebook collection and exports "
        "a printable PDF worksheet — one failure case per page."
    )

    architecture_label = "Tablebook RAG"

    if st.button("🚀 Run all test questions & generate PDF"):
        results = []
        progress = st.progress(0)
        status = st.empty()

        for i, q in enumerate(TEST_QUESTIONS):
            status.write(f"Running {q['id']}: {q['question'][:60]}…")
            chunks, metas = retrieve_chunks(q["question"])
            context = "\n---\n".join(chunks) if chunks else ""
            answer = ask_llm(context, q["question"])
            sources = [m.get("source", "?") for m in metas]
            results.append({
                "id":       q["id"],
                "question": q["question"],
                "category": q["category"],
                "chunks":   chunks,
                "sources":  sources,
                "answer":   answer,
            })
            progress.progress((i + 1) / len(TEST_QUESTIONS))

        status.write("✅ All questions run — building PDF…")

        try:
            pdf_bytes = build_failure_pdf(results, architecture_label)
            filename = (f"failure_worksheet_tablebook_rag_"
                        f"{datetime.now().strftime('%Y%m%d')}.pdf")
            st.success(f"PDF ready! {len(results)} cases exported.")
            st.download_button(
                label="⬇️ Download Worksheet PDF",
                data=pdf_bytes,
                file_name=filename,
                mime="application/pdf",
            )
        except Exception as e:
            st.error(f"PDF generation failed: {e}")

    st.markdown("---")
    st.markdown("**Failure category codes:**")
    cols = st.columns(3)
    codes = [
        ("R", "Retrieval Miss",    "Wrong chunks fetched"),
        ("C", "Chunking Break",    "Answer split across chunk boundary"),
        ("T", "Table/Format Loss", "Table structure destroyed in parsing"),
        ("I", "Image Blind Spot",  "Content was a diagram or image"),
        ("H", "Hallucination",     "LLM invented facts not in context"),
        ("L", "Layout Error",      "Multi-column or footnote parsed wrong"),
    ]
    for i, (code, name, desc) in enumerate(codes):
        with cols[i % 3]:
            st.markdown(f"**{code} — {name}**  \n{desc}")